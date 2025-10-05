from pptx import Presentation

def extract_pptx(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide_num, slide in enumerate(prs.slides, start=1):
        text += f"\n[SLIDE {slide_num}]\n"
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text
    